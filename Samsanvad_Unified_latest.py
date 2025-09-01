from flask import session
from flask import url_for
from langchain.vectorstores import FAISS
from langchain.embeddings import OpenAIEmbeddings
import traceback
import pandas as pd
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
import faiss
from langchain.docstore.document import Document
import base64
from flask import jsonify, send_file
import uuid, os
from flask import Flask, render_template, request, redirect
import json
from openai import OpenAI
from werkzeug.utils import secure_filename
from datetime import datetime
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
import secrets

client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))  # or set directly

# Folder setup
GLOBAL_UPLOAD_FOLDER = "global_uploads"  #  save all documents in this folder
MULTI_UPLOAD_FOLDER = "index_Individual_uploads"  # saves all individual documents in this folder.
GLOBAL_INDEX_DIR = "index_Global_Vector"
QUERY_UPLOAD_FOLDER = "query_upload"
USERS_FILE ="users.json"
os.makedirs(GLOBAL_UPLOAD_FOLDER, exist_ok=True)
os.makedirs(MULTI_UPLOAD_FOLDER, exist_ok=True)
os.makedirs(GLOBAL_INDEX_DIR, exist_ok=True)

app = Flask(__name__)
app.secret_key = secrets.token_hex(32)
@app.route('/')
def home():
    return redirect(url_for('login'))

@app.route('/login', methods=['GET', 'POST'])
def login():
    error = None
    success = None

    if request.method == 'POST':
        email = request.form['email'].strip().lower()
        username = request.form['username'].strip()
        password = request.form['password'].strip()

        if not os.path.exists(USERS_FILE):
            error = "⚠️ User database not found."
        else:
            with open(USERS_FILE, 'r') as f:
                users = json.load(f)

            user_record = users.get(email)
            if user_record:
                username_match = user_record.get('username', '').strip().lower() == username.lower()
                password_match = user_record.get('password', '') == password
                if username_match and password_match:
                    session['username'] = username
                    return redirect(url_for("main_page"))
                else:
                    error = "❌ Invalid email, username, or password."
            else:
                error = "❌ Invalid email, username, or password."

    return render_template('login.html', error=error, success=success, form_action=url_for('login'))

@app.route('/register', methods=['GET', 'POST'])
def register():
    error = None
    success = None

    if request.method == 'POST':
        email = request.form['email'].strip().lower()
        username = request.form['username'].strip()
        password = request.form['password'].strip()

        users = {}
        if os.path.exists(USERS_FILE):
            with open(USERS_FILE, 'r') as f:
                users = json.load(f)

        if email in users:
            error = "⚠️ Email already registered."
        else:
            users[email] = {
                "username": username,
                "password": password
            }
            with open(USERS_FILE, 'w') as f:
                json.dump(users, f, indent=2)
            success = "✅ Registration successful. You can now sign in."

    return render_template('login.html', error=error, success=success, form_action=url_for('register'))

@app.route('/main_page')
def main_page():
    username = session.get('username', 'User')
    return render_template('ChatSideBar5.html', username=username)

UPLOAD_FOLDER = 'static'
LOGO_FILENAME = 'logo6.jpg'  # Default logo

CONTEXT_FILE = "user_context.json"
DEFAULT_CONTEXT = {
    "role": "guest",
    "context": "default"
}
user_context = DEFAULT_CONTEXT.copy()

@app.route('/save_context', methods=['POST'])
def save_context():
    global user_context
    data = request.get_json()
    print("🔔 Received save_context:", data)

    role = data.get('role') or DEFAULT_CONTEXT["role"]
    context = data.get('context') or DEFAULT_CONTEXT["context"]

    user_context["role"] = role
    user_context["context"] = context

    try:
        with open(CONTEXT_FILE, 'w') as f:
            json.dump(user_context, f, indent=2)
        return jsonify({"status": "saved", "data": user_context})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/load_context', methods=['GET'])
def load_context():
    try:
        with open(CONTEXT_FILE) as f:
            user_context = json.load(f)
        return jsonify(user_context)
    except Exception:
        return jsonify(DEFAULT_CONTEXT)

@app.route('/get_dashboard', methods=['GET'])
def get_dashboard():
    role = request.headers.get('X-User-Role') or DEFAULT_CONTEXT["role"]
    context = request.headers.get('X-Context') or DEFAULT_CONTEXT["context"]

    print(f"🔍 get_dashboard: role={role}, context={context}")

    DASHBOARD_MODULES = [
        {"id": "analytics", "label": "📊 Analytics", "roles": ["admin", "analyst"]},
        {"id": "editor", "label": "✏️ Editor", "roles": ["admin", "editor"]},
        {"id": "viewer", "label": "👁️ Viewer", "roles": ["admin", "viewer", "guest"]},
    ]

    filtered_modules = [
        module for module in DASHBOARD_MODULES
        if role in module.get("roles", [])
    ]

    return jsonify({
        "role": role,
        "context": context,
        "modules": filtered_modules
    })

@app.route("/upload_logo", methods=["POST"])
def upload_logo():
    file = request.files.get("logo")
    if file and file.filename:
        filepath = os.path.join(UPLOAD_FOLDER, LOGO_FILENAME)
        file.save(filepath)
        return "Logo uploaded successfully", 200
    return "No file received", 400

 # Unified path
def load_excel_as_documents(file_path):
    df = pd.read_excel(file_path, engine="openpyxl")

    documents = [
        Document(
            page_content=", ".join([
                f"{col}: {row[col]}" for col in df.columns if pd.notnull(row[col])
            ]),
            metadata={"row_index": idx}
        )
        for idx, row in df.iterrows()
    ]
    return documents


@app.route("/upload", methods=["POST"])
def upload():
    file = request.files["file"]

    #vector_mode = request.args.get("mode", "single")
    vector_mode = request.form.get("vectorMode","single")
    #print(vector_mode)

    filename = os.path.splitext(file.filename)[0]

    if vector_mode == "multi":
        upload_folder = MULTI_UPLOAD_FOLDER
        index_dir = os.path.join(upload_folder, f"{filename}_index")
    else:
        upload_folder = GLOBAL_UPLOAD_FOLDER
        index_dir = GLOBAL_INDEX_DIR

    os.makedirs(upload_folder, exist_ok=True)
    filepath = os.path.join(upload_folder, file.filename)
    file.save(filepath)

    append_to_index(filepath, index_dir)
    return jsonify({"message": f"{file.filename} uploaded successfully in '{vector_mode}' mode."})

@app.route("/uploadquery", methods=["POST"])
def uploadquery():
    file = request.files.get("file")
    if not file:
        return jsonify({"error": "❌ No file provided."}), 400

    filename = file.filename
    upload_folder = QUERY_UPLOAD_FOLDER  # Always use this folder
    os.makedirs(upload_folder, exist_ok=True)

    filepath = os.path.join(upload_folder, filename)
    file.save(filepath)

    print(f"✅ Query document saved at: {filepath}")

    return jsonify({
        "message": f"✅ File '{filename}' uploaded successfully to 'query_upload'. No index created."
    })


@app.route("/listqueries", methods=["GET"])
def list_uploaded_queries():
    upload_folder = QUERY_UPLOAD_FOLDER
    os.makedirs(upload_folder, exist_ok=True)

    files = [
        f for f in os.listdir(upload_folder)
        if os.path.isfile(os.path.join(upload_folder, f))
    ]

    return jsonify({"query_documents": files})


# ---------- INDEXING ----------
def append_to_index(file_path, index_dir):
    extension = os.path.splitext(file_path)[1].lower()

    if extension == '.pdf':
        from langchain.document_loaders import PyPDFLoader
        loader = PyPDFLoader(file_path)
        new_docs = loader.load()

    elif extension == '.txt':
        from langchain.document_loaders import TextLoader
        loader = TextLoader(file_path, encoding="utf-8")
        new_docs = loader.load()

    elif extension in ['.doc', '.docx']:
        from langchain.document_loaders import UnstructuredWordDocumentLoader
        loader = UnstructuredWordDocumentLoader(file_path)
        new_docs = loader.load()

    elif extension in ['.xls', '.xlsx']:
        new_docs = load_excel_as_documents(file_path)

    elif extension == '.xml':
        from langchain.document_loaders import UnstructuredXMLLoader
        loader = UnstructuredXMLLoader(file_path)
        new_docs = loader.load()

    else:
        raise ValueError(f"Unsupported file type: {extension}")

    embeddings = OpenAIEmbeddings()

    try:
        faiss_path = os.path.join(index_dir, "index.faiss")
        if os.path.exists(faiss_path) and os.path.getsize(faiss_path) > 0:
            db = FAISS.load_local(index_dir, embeddings)
            db.add_documents(new_docs)
            print(f"✅ Appended to index: {index_dir}")
        else:
            raise FileNotFoundError("Index missing or empty.")
    except Exception as e:
        print(f"⚠️ Load failed: {e}. Creating new index...")
        db = FAISS.from_documents(new_docs, embeddings)

    os.makedirs(index_dir, exist_ok=True)
    db.save_local(index_dir)

    try:
        import faiss
        reloaded = faiss.read_index(os.path.join(index_dir, "index.faiss"))
        print("✅ FAISS reload successful.")
    except Exception as e:
        print(f"❌ FAISS reload failed: {e}")

@app.route("/files", methods=["GET"])
def list_files():
    vector_mode = request.args.get("mode", "single")

    if vector_mode == "multi":
        files = [
            name.replace("_index", "")
            for name in os.listdir(MULTI_UPLOAD_FOLDER)
            if os.path.isdir(os.path.join(MULTI_UPLOAD_FOLDER, name)) and name.endswith("_index")
        ]
    else:
        files = [
            name for name in os.listdir(GLOBAL_UPLOAD_FOLDER)
            if os.path.isfile(os.path.join(GLOBAL_UPLOAD_FOLDER, name))
        ]
    #print(files)
    return jsonify({"files": files})

@app.route("/index_url", methods=["POST"])
def index_url():
    url = request.json.get("url")
    vector_mode = request.json.get("vectorMode", "single")
    domain_name = urlparse(url).netloc.replace(".", "_")  # Use domain name as identifier
    if vector_mode == "multi":
        upload_folder = MULTI_UPLOAD_FOLDER
        index_dir = os.path.join(upload_folder, f"{domain_name}_index")
    else:
        upload_folder = GLOBAL_UPLOAD_FOLDER
        index_dir = GLOBAL_INDEX_DIR

    os.makedirs(upload_folder, exist_ok=True)
    os.makedirs(index_dir, exist_ok=True)

    success = crawl_and_index_site(url, index_dir)
    if success:
        return jsonify({"message": f"Website indexed successfully from '{url}'."})
    return jsonify({"message": "Failed to index website."}), 500

def crawl_and_index_site1(base_url, index_dir, max_pages=50):
    print('Crowl')
    visited = set()
    to_visit = [base_url]
    new_docs = []
    embeddings = OpenAIEmbeddings()

    while to_visit and len(visited) < max_pages:
        url = to_visit.pop(0)
        if url in visited or urlparse(url).netloc != urlparse(base_url).netloc:
            continue

        try:
            res = requests.get(url)
            soup = BeautifulSoup(res.text, "html.parser")
            text = soup.get_text(separator=" ", strip=True)

            metadata = { "source": url }
            doc = Document(page_content=text, metadata=metadata)
            new_docs.append(doc)
            visited.add(url)

            for link in soup.find_all("a", href=True):
                full_link = urljoin(url, link['href'])
                if full_link.startswith(base_url) and full_link not in visited:
                    to_visit.append(full_link)

        except Exception as e:
            print(f"❌ Error at {url}: {e}")

    try:
        faiss_path = os.path.join(index_dir, "index.faiss")
        if os.path.exists(faiss_path) and os.path.getsize(faiss_path) > 0:
            db = FAISS.load_local(index_dir, embeddings,allow_dangerous_deserialization=True)
            db.add_documents(new_docs)
            print(f"✅ Appended website content to index: {index_dir}")
        else:
            raise FileNotFoundError("Index missing or empty.")
    except Exception as e:
        print(f"⚠️ Reload failed: {e}. Creating new index...")
        db = FAISS.from_documents(new_docs, embeddings)

    db.save_local(index_dir)
    save_crawled_content_txt(new_docs, index_dir)

    print(f"📄 Total pages crawled: {len(visited)}")

    try:
        reloaded = faiss.read_index(faiss_path)
        print("✅ FAISS reload successful.")
        return True
    except Exception as e:
        print(f"❌ Indexing error: {type(e).__name__} - {e}")
        traceback.print_exc()
        print("⚠️ Reload failed. Creating new index...")
        db = FAISS.from_documents(new_docs, embeddings)

def save_crawled_content_txt(docs, index_dir):
    output_path = os.path.join(index_dir, "crawled_pages.txt")
    with open(output_path, "w", encoding="utf-8") as f:
        for doc in docs:
            f.write(f"🔗 Source: {doc.metadata['source']}\n")
            f.write(doc.page_content + "\n\n" + "="*80 + "\n\n")
    print(f"📁 Saved crawled content to: {output_path}")

def chunk_docs(docs, max_tokens=250000):
    batches = []
    current_batch = []
    current_tokens = 0

    for doc in docs:
        tokens = len(doc.page_content.split())
        if current_tokens + tokens > max_tokens:
            batches.append(current_batch)
            current_batch = []
            current_tokens = 0
        current_batch.append(doc)
        current_tokens += tokens

    if current_batch:
        batches.append(current_batch)
    return batches

def crawl_and_index_site(base_url, index_dir, max_pages=20, output_dir="global_uploads"):
    print("🐛 Crowl")
    visited = set()
    to_visit = [base_url]
    new_docs = []
    embeddings = OpenAIEmbeddings()

    while to_visit and len(visited) < max_pages:
        url = to_visit.pop(0)
        if url in visited or urlparse(url).netloc != urlparse(base_url).netloc:
            continue

        try:
            res = requests.get(url)
            soup = BeautifulSoup(res.text, "html.parser")
            text = soup.get_text(separator=" ", strip=True)

            if len(text.split()) > 50000:
                print(f"⚠️ Skipping oversized page: {url}")
                continue

            metadata = { "source": url }
            doc = Document(page_content=text, metadata=metadata)
            new_docs.append(doc)
            visited.add(url)

            for link in soup.find_all("a", href=True):
                full_link = urljoin(url, link['href'])
                if full_link.startswith(base_url) and full_link not in visited:
                    to_visit.append(full_link)
        except Exception as e:
            print(f"❌ Error at {url}: {e}")

    db = None
    faiss_path = os.path.join(index_dir, "index.faiss")

    try:
        if os.path.exists(faiss_path) and os.path.getsize(faiss_path) > 0:
            db = FAISS.load_local(index_dir, embeddings, allow_dangerous_deserialization=True)
            print(f"📦 Loaded existing index.")
        else:
            raise FileNotFoundError("Index missing or empty.")
    except Exception as e:
        print(f"❌ Indexing error: {type(e).__name__} - {e}")
        traceback.print_exc()
        print("⚠️ Reload failed. Creating new index...")

    for batch in chunk_docs(new_docs):
        if db:
            db.add_documents(batch)
        else:
            db = FAISS.from_documents(batch, embeddings)

    db.save_local(index_dir)
    save_crawled_content_txt(new_docs, output_dir, filename="samsan_crawled.txt")

    print(f"📄 Total pages crawled: {len(visited)}")

    try:
        reloaded = faiss.read_index(faiss_path)
        print("✅ FAISS reload successful.")
        return True
    except Exception as e:
        print(f"❌ Final FAISS reload failed: {e}")
        return False

latest_chat_response = None


@app.route("/chat", methods=["POST"])
def chat():
    data = request.get_json()
    query = data.get("query", "")
    vector_mode = data.get("vectorMode")
    print('chat_vector_mode')
    filename = data.get("filename", None)

    embeddings = OpenAIEmbeddings()

    # Load correct index based on mode
    if vector_mode == "multi" and filename:
        index_dir = os.path.join(MULTI_UPLOAD_FOLDER, f"{filename}_index")
    else:
        index_dir = GLOBAL_INDEX_DIR
    print(index_dir)
    try:
        db = FAISS.load_local(index_dir, embeddings, allow_dangerous_deserialization=True)
    except Exception as e:
        return jsonify({"response": f"⚠️ Could not load FAISS index: {str(e)}"})

    # Retrieve relevant docs and construct context
    similar_docs = db.similarity_search(query, k=4)
    context = "\n\n".join([doc.page_content for doc in similar_docs])

    try:
        messages = [
            {"role": "system", "content": "Use the context to answer clearly and helpfully."},
            {"role": "user", "content": f"Context:\n{context}\n\nUser asked: {query}"}
        ]

        response = client.chat.completions.create(
            model="gpt-4",
            messages=messages
        )

        final_response = response.choices[0].message.content.strip()
        latest_chat_response = final_response  # 🧠 Save response for PPT
        # Optional PPT generation block
        # slide_data = {"topic": query, "response": final_response}
        # ppt_path = create_ppt(slide_data)

        # return jsonify({
        #     "response": final_response,
        #     "ppt_created": True,
        #     "ppt_file": ppt_path
        # })

        return jsonify({"response": final_response})

    except Exception as e:
        error_message = str(e)
        if "401" in error_message or "Invalid" in error_message:
            return jsonify({"response": "❌ OpenAI API key is missing or invalid."})
        return jsonify({"response": f"⚠️ Chat model error: {error_message}"})

@app.route("/docquery", methods=["POST"])
def docquery():
    print('docquery')
    data = request.get_json()
    query_filename = data.get("queryfilename ", None)
    print('filename', query_filename)  # Excel query file
    knowledge_filename = data.get("knowledgeFilename", None)  # Indexed document
    print('knowledge_filename', knowledge_filename)
    vector_mode = data.get("vectorMode", "multi")
    print('vectorMode', vector_mode)

    if not query_filename or not knowledge_filename:
        return jsonify({"response": "⚠️ Missing query or knowledge filename."}), 400

    # Load FAISS index
    try:
        embeddings = OpenAIEmbeddings()
        base_name = os.path.splitext(knowledge_filename)[0]
        index_dir = (
            os.path.join(MULTI_UPLOAD_FOLDER, f"{base_name}_index")
            if vector_mode == "multi"
            else GLOBAL_INDEX_DIR
        )
        db = FAISS.load_local(index_dir, embeddings, allow_dangerous_deserialization=True)
    except Exception as e:
        return jsonify({"response": f"⚠️ Could not load FAISS index: {str(e)}"}), 500

    query_path = os.path.join(QUERY_UPLOAD_FOLDER, query_filename)
    file_ext = os.path.splitext(query_filename)[-1].lower()

    queries = []
    cell_map = []

    try:
        if file_ext == ".xlsx":
            wb = load_workbook(query_path)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for row in ws.iter_rows(min_row=2, max_col=ws.max_column):
                    for cell in row:
                        query_text = str(cell.value).strip() if cell.value else ""
                        if query_text:
                            queries.append(query_text)
                            cell_map.append(("excel", sheet_name, cell.row, cell.column))
        elif file_ext == ".csv":
            import pandas as pd
            df = pd.read_csv(query_path)
            for row_idx, row in df.iterrows():
                for col_idx, value in enumerate(row):
                    query_text = str(value).strip() if pd.notna(value) else ""
                    if query_text:
                        queries.append(query_text)
                        cell_map.append(("csv", row_idx, col_idx))
        else:
            return jsonify({"response": "❌ Unsupported file type."}), 400
    except Exception as e:
        return jsonify({"response": f"⚠️ Failed to read query file: {str(e)}"}), 500

    if not queries:
        return jsonify({"response": "⚠️ No queries found in the file."}), 400

    # Retrieve shared context
    try:
        similar_docs = db.similarity_search("\n".join(queries), k=6)
        indexed_context = "\n\n".join([doc.page_content for doc in similar_docs])
    except Exception as e:
        return jsonify({"response": f"⚠️ Similarity search failed: {str(e)}"}), 500

    # Construct prompt
    messages = [
        {
            "role": "system",
            "content": (
                "You are a precision assistant. For each query below, return only a numeric answer if available. "
                "If no numeric answer exists, return 'NA'. Keep fallback responses under 5 words. "
                "Do not include explanations, symbols, or units. Respond as a numbered list matching the query order."
            )
        },
        {
            "role": "user",
            "content": (
                    f"Context:\n{indexed_context}\n\n"
                    f"Queries:\n" + "\n".join([f"{i + 1}. {q}" for i, q in enumerate(queries)])
            )
        }
    ]

    # Call GPT-4
    try:
        response = client.chat.completions.create(model="gpt-4", messages=messages)
        raw_output = response.choices[0].message.content.strip()
        answers = [line.split(". ", 1)[-1].strip() for line in raw_output.split("\n") if line.strip()]
    except Exception as e:
        return jsonify({"response": f"⚠️ Chat model error: {str(e)}"}), 500

    # Write responses
    try:
        output_filename = f"processed_{query_filename}"
        output_path = os.path.join(QUERY_UPLOAD_FOLDER, output_filename)

        if file_ext == ".xlsx":
            for (source, sheet_name, row, col), answer in zip(cell_map, answers):
                wb[sheet_name].cell(row=row, column=col + 1).value = answer
            wb.save(output_path)
        elif file_ext == ".csv":
            for (source, row_idx, col_idx), answer in zip(cell_map, answers):
                df.iat[row_idx, col_idx + 1] = answer
            df.to_csv(output_path, index=False)

        return jsonify({"response": f"✅ Responses written to '{output_filename}' successfully."})
    except Exception as e:
        return jsonify({"response": f"⚠️ Failed to write responses: {str(e)}"}), 500

def generate_filename(topic):
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_topic = sanitize_topic(topic)[:20]
    return f"{safe_topic}_{ts}.pptx"


@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    print("📥 /generate_ppt reached")
    try:
        data = request.get_json(force=True)
        print("📦 Received data:", data)

        slide_count = data.get("slide_count")
        try:
            slide_count = int(slide_count)
            if slide_count < 1:
                raise ValueError
        except ValueError:
            print("❌ Invalid slide count")
            return jsonify({
                "status": "error",
                "message": "❌ Invalid slide count. Must be a positive integer."
            }), 400

        print("✅ Slide count confirmed:", slide_count)

        output_dir = "C:/PythonProject/static/ppt"
        os.makedirs(output_dir, exist_ok=True)

        chat_json = data.get("chat_json")
        slides_json = data.get("slides_json", [])
        text = data.get("text", "").strip()
        topic = data.get("topic", "presentation")

        filename = None
        output_path = None

        if chat_json:
            print("🧠 Using chat_json")
            filename = generate_slides_from_chat(chat_json, slide_count, output_dir)
            output_path = os.path.join(output_dir, filename)

        elif isinstance(slides_json, list) and slides_json:
            #print("📊 Using slides_json")
            formatted_text = json_to_paragraphs(slides_json)
            slides = parse_generated_content(formatted_text, slide_count)
            filename = generate_filename(topic)
            output_path = os.path.join(output_dir, filename)
            create_presentation(slides, output_path)

        elif text:
            #print("📄 Using raw text")
            slides = parse_generated_content(text, slide_count)
            filename = generate_filename(topic)
            output_path = os.path.join(output_dir, filename)
            create_presentation(slides, output_path)

        else:
            print("❌ No valid input")
            return jsonify({
                "status": "error",
                "message": "❌ No valid input provided for slide generation."
            }), 400

        # ✅ Confirm file exists before returning
        if not os.path.exists(output_path):
            print("⚠️ File not found after creation:", output_path)
            return jsonify({
                "status": "error",
                "message": "⚠️ File was not saved correctly."
            }), 500

        #print("✅ File saved:", output_path)
        #print("📤 Returning success response")

        return jsonify({
            "status": "success",
            "message": f"✅ Presentation saved to: {output_path}",
            "path": output_path,
            "download_url": f"/download_ppt?file={filename}"
        }), 200

    except Exception as e:
        print("🔥 Exception caught:", str(e))
        return jsonify({
            "status": "error",
            "message": f"⚠️ Slide generation failed: {str(e)}"
        }), 500

def generate_slides_from_chat(chat_json, slide_count=5, output_dir="C:/PythonProject/static/ppt"):
    response_text = get_slide_response_from_llm(chat_json, slide_count)

    # response_text = """
    # Slide 1: Introduction to Quantum Computing
    # - Quantum computing uses qubits instead of classical bits.
    # - It leverages superposition and entanglement.
    # - Promises exponential speedup for certain problems.
    #
    # Slide 2: Applications of Quantum Computing
    # - Drug discovery and molecular simulation.
    # - Optimization in logistics and finance.
    # - Cryptography and secure communication.
    #
    # Slide 3: Challenges in Quantum Computing
    # - Qubit decoherence and error rates.
    # - Scalability of quantum hardware.
    # - Need for specialized algorithms.
    #
    # Slide 4: Quantum vs Classical Computing
    # - Classical bits are binary; qubits can be in multiple states.
    # - Quantum computers solve certain problems faster.
    # - Classical systems still dominate general-purpose computing.
    #
    # Slide 5: Future of Quantum Computing
    # - Rapid advancements in hardware and theory.
    # - Potential to revolutionize industries.
    # - Still in early stages of practical deployment.
    # """

    slides = parse_generated_content(response_text, slide_count)

    uid = uuid.uuid4().hex[:6]
    filename = f"presentation_{uid}.pptx"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, filename)

    create_presentation(slides, output_path)
    return filename

def parse_generated_content_Title_Bullet(response, slide_count=2):
    def normalize_bullets(lines):
        return [line.strip("-• ").strip() for line in lines if line.strip()]

    def deduplicate_preserve_order(items):
        seen = set()
        return [x for x in items if not (x in seen or seen.add(x))]

    slides = []

    # Try structured parsing first
    chunks = re.split(r"(?:Slide\s*\d+:)", response)
    titles = re.findall(r"Slide\s*\d+:\s*(.*?)\n", response)

    if len(chunks) > 1 and titles:
        for idx, chunk in enumerate(chunks[1:]):
            title = titles[idx].strip()
            bullets = normalize_bullets(chunk.strip().split("\n"))
            bullets = deduplicate_preserve_order(bullets)
            slides.append({"title": title, "bullets": bullets})
            if len(slides) >= slide_count:
                break

    # Fallback: split by sentence if no line breaks
    if len(slides) < slide_count:
        lines = normalize_bullets(response.strip().split("\n"))
        if len(lines) <= 1:
            lines = re.split(r'(?<=[.!?])\s+', response.strip())
            lines = normalize_bullets(lines)

        lines = deduplicate_preserve_order(lines)

        total_lines = len(lines)
        per_slide = max(1, total_lines // slide_count)
        remainder = total_lines % slide_count
        start = 0
        for i in range(slide_count):
            extra = 1 if i < remainder else 0
            end = start + per_slide + extra
            chunk = lines[start:end]
            if not chunk:
                chunk = [f"Placeholder content {i+1}"]
            slides.append({
                "title": f"Slide {len(slides)+1}",
                "bullets": chunk
            })
            start = end

    return slides

def parse_generated_content(response, slide_count=2):
    def normalize_bullets(lines):
        return [line.strip("-• ").strip() for line in lines if line.strip()]

    slides = []

    # Try structured parsing first
    chunks = re.split(r"(?:Slide\s*\d+:)", response)
    titles = re.findall(r"Slide\s*\d+:\s*(.*?)\n", response)

    if len(chunks) > 1 and titles:
        for idx, chunk in enumerate(chunks[1:]):
            title = titles[idx].strip()
            lines = chunk.strip().split("\n")
            bullets = normalize_bullets(lines)

            # Remove title if it's repeated as first bullet
            if bullets and bullets[0].lower() == title.lower():
                bullets = bullets[1:]

            slides.append({"title": title, "bullets": bullets})
            if len(slides) >= slide_count:
                break

    # Fallback: split by sentence if no line breaks
    if len(slides) < slide_count:
        lines = normalize_bullets(response.strip().split("\n"))
        if len(lines) <= 1:
            lines = re.split(r'(?<=[.!?])\s+', response.strip())
            lines = normalize_bullets(lines)

        total_lines = len(lines)
        per_slide = max(1, total_lines // slide_count)
        remainder = total_lines % slide_count
        start = 0
        for i in range(slide_count):
            extra = 1 if i < remainder else 0
            end = start + per_slide + extra
            chunk = lines[start:end]
            if not chunk:
                chunk = [f"Placeholder content {i+1}"]
            slides.append({
                "title": f"Slide {len(slides)+1}",
                "bullets": chunk
            })
            start = end

    return slides

def create_presentation(slides, output_file):
    prs = Presentation()
    base_image_path = "static/images/"

    # Preload available images
    available_images = [
        f for f in os.listdir(base_image_path)
        if f.lower().endswith((".png", ".jpg", ".jpeg"))
    ]
    print(f"🖼️ Found {len(available_images)} images in {base_image_path}")

    for i, slide_data in enumerate(slides):
        title = slide_data.get("title", f"Slide {i+1}")
        bullets = slide_data.get("bullets", [])
        image_filename = slide_data.get("image_filename")  # Optional
        image_align = slide_data.get("image_align", "right")
        footer_label = slide_data.get("footer_label", "QuantumDeck")
        footer_text = f"Slide {i+1} • {footer_label}"

        # If no image specified, pick one randomly
        if not image_filename and available_images:
            image_filename = random.choice(available_images)
            print(f"🎲 Slide {i+1}: Random image selected → {image_filename}")

        image_path = os.path.join(base_image_path, image_filename) if image_filename else None

        generate_slide(
            prs,
            title_text=title,
            bullets=bullets,
            image_path=image_path,
            image_align=image_align,
            footer_text=footer_text,
            slide_index=i
        )

    prs.save(output_file)

def generate_slide(prs, title_text, bullets, image_path=None, image_align="right", footer_text=None,
                   slide_index=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Content
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT

    # ✅ Optional Image
    if image_path:
        img_width = Inches(2.5)
        img_height = Inches(2.5)
        positions = {
            "right": (Inches(7.0), Inches(2.0)),
            "left": (Inches(0.5), Inches(2.0)),
            "top": (Inches(3.5), Inches(0.5)),
            "bottom": (Inches(3.5), Inches(5.0))
        }
        left, top = positions.get(image_align, positions["right"])
        slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)

    # ✅ Footer Text
    from datetime import datetime

    # Footer (optional)
    if footer_text:
        left = Inches(0.5)
        top = Inches(6.8)
        width = Inches(9)
        height = Inches(0.5)
        footer_box = slide.shapes.add_textbox(left, top, width, height)
        tf = footer_box.text_frame
        tf.clear()

        # Compose footer: Slide number • Date
        slide_number = f"Slide {slide_index + 1}"
        current_date = datetime.now().strftime("%b %d, %Y")  # e.g. "Aug 18, 2025"
        full_footer = f"{slide_number} • {current_date}"

        p = tf.paragraphs[0]
        p.text = full_footer
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = 1  # Center

    return slide


def json_to_paragraph_string(slides_json):
    result = ""
    for item in slides_json:
        result += f"{item.get('title', 'Untitled Slide')}:\n"
        for bullet in item.get("bullets", []):
            result += f"- {bullet}\n"
        result += "\n"
    return result.strip()

def auto_chunk_text(response_text, desired_slide_count=5):
    lines = [line.strip() for line in response_text.split("\n") if line.strip()]
    chunk_size = max(1, len(lines) // desired_slide_count)

    chunks = []
    for i in range(0, len(lines), chunk_size):
        chunk = lines[i:i + chunk_size]
        title = f"Slide {len(chunks)+1}"
        bullets = chunk
        chunks.append({"title": title, "bullets": bullets})
    return chunks

# Build the presentation from parsed slides
# Convert PPT to base64 (optional use)
def ppt_to_base64(filepath):
    with open(filepath, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")

def json_to_paragraphs(slides_json):
    formatted = ""
    for item in slides_json:
        formatted += f"{item['title']}:\n"
        for b in item.get("bullets", []):
            formatted += f"- {b}\n"
        formatted += "\n"
    return formatted.strip()


@app.route('/download_ppt')
def download_ppt():
    file_param = request.args.get("file", "")
    filename = secure_filename(file_param)
    filepath = os.path.join("C:/PythonProject/static/ppt", filename)

    if not os.path.exists(filepath):
        return "❌ File not found.", 404

    return send_file(filepath, as_attachment=True)

def sanitize_topic(topic):
    return topic.strip().replace(" ", "_").lower()

def normalize_bullets(lines):
    bullets = []
    for line in lines:
        bullet = re.sub(r"^\s*[-•*•\d.]+", "", line).strip()
        if bullet:
            bullets.append(bullet)
    return bullets

def build_slide_generation_prompt(chat_json, slide_count=5):
    prompt = f"""
        Reformat the following chatbot JSON output into clearly structured presentation content.
        
        Include approximately {slide_count} slides. Each slide should follow this format:
        Slide <number>: <Slide Title>
        - Bullet point 1
        - Bullet point 2
        - Bullet point 3
        
        Here's the JSON:
        {json.dumps(chat_json, indent=2)}
        """
    return prompt.strip()

def get_slide_response_from_llm(chat_json, slide_count=5):
    prompt = build_slide_generation_prompt(chat_json, slide_count)

    messages = [
        {"role": "system", "content": "You are a helpful assistant that formats responses into presentations."},
        {"role": "user", "content": prompt}
    ]

    response = client.chat.completions.create(
        model="gpt-4",
        messages=messages
    )

    return response.choices[0].message.content.strip()

if __name__ == "__main__":
    app.run(debug=True)